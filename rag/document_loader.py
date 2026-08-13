"""
DocumentLoader — pure pip, no system dependencies.
Scanned PDFs → Groq vision model (same API key).
"""
import re, json, csv, io, xml.etree.ElementTree as ET
from pathlib import Path
from collections import Counter
from typing import List, Dict, Any
from utils.logger import setup_logger

logger = setup_logger(__name__)


class DocumentLoader:
    def __init__(self, config=None, groq_client=None):
        self.config      = config
        self.groq_client = groq_client

    def load(self, filepath: str) -> List[Dict[str, Any]]:
        ext      = Path(filepath).suffix.lower().lstrip(".")
        filename = Path(filepath).name
        loaders  = {
            "pdf":self._pdf,"txt":self._text,"log":self._text,"md":self._text,
            "markdown":self._text,"docx":self._docx,"doc":self._docx,
            "csv":self._csv,"xlsx":self._xlsx,"xls":self._xlsx,
            "pptx":self._pptx,"ppt":self._pptx,"json":self._json,
            "xml":self._xml,"html":self._html,"htm":self._html,
            "png":self._image,"jpg":self._image,"jpeg":self._image,
            "bmp":self._image,"tiff":self._image,"webp":self._image,
        }
        fn = loaders.get(ext)
        if not fn: logger.warning(f"No loader: {ext}"); return []
        try: pages = fn(filepath)
        except Exception as e: logger.error(f"Load failed {filename}: {e}"); return []
        if not pages: logger.warning(f"No content: {filename}"); return []
        for p in pages:
            p.setdefault("metadata",{})
            p["metadata"]["source"]    = filename
            p["metadata"]["file_type"] = ext
        pages = self._rm_noise(pages)
        if self.config and getattr(self.config,"PII_MASKING_ENABLED",True):
            pages = self._mask_pii(pages)
        logger.info(f"Loaded {len(pages)} sections from {filename}")
        return pages

    # ── PDF ─────────────────────────────────────────────
    def _pdf(self, fp):
        if self._encrypted(fp):
            return [{"content":"[Password-protected PDF — please unlock first]","metadata":{}}]
        out = []
        try:
            import pdfplumber
            with pdfplumber.open(fp) as pdf:
                for i,page in enumerate(pdf.pages):
                    pn   = i+1
                    text = self._multicol(page)
                    tbl  = self._pdf_tables(page)
                    if tbl: text = (text or "") + "\n\nTABLES:\n" + tbl
                    if not text or len(text.strip()) < 20:
                        logger.info(f"Page {pn}: no text — trying vision")
                        text = self._vision_page(page, pn)
                    if text and text.strip():
                        out.append({"content":text.strip(),"metadata":{"page":pn}})
                    # Embedded images → vision
                    if self.groq_client and getattr(self.config,"VISION_ENABLED",True):
                        for d in self._page_images(page, pn): out.append(d)
        except Exception as e:
            logger.error(f"pdfplumber error: {e}")
        if not out:
            out=[{"content":"[Scanned PDF — content could not be fully extracted. Try a text-based PDF.]","metadata":{"page":1}}]
        return out

    def _encrypted(self, fp):
        try:
            import pdfplumber
            with pdfplumber.open(fp) as pdf: _ = pdf.pages[0].extract_text()
            return False
        except Exception as e:
            return "encrypt" in str(e).lower() or "password" in str(e).lower()

    def _multicol(self, page):
        try:
            words = page.extract_words(x_tolerance=3, y_tolerance=3)
            if not words: return page.extract_text() or ""
            mid   = page.width / 2
            left  = [w for w in words if float(w["x0"]) < mid]
            right = [w for w in words if float(w["x0"]) >= mid]
            if len(right) > max(len(left)*0.3, 10):
                left.sort( key=lambda w:(round(float(w["top"])/5)*5,float(w["x0"])))
                right.sort(key=lambda w:(round(float(w["top"])/5)*5,float(w["x0"])))
                return " ".join(w["text"] for w in left) + "\n\n" + " ".join(w["text"] for w in right)
            return page.extract_text() or ""
        except: return page.extract_text() or ""

    def _pdf_tables(self, page):
        try:
            tables = page.extract_tables()
            if not tables: return ""
            rows = []
            for t in tables:
                for row in t:
                    cells = [str(c).strip() if c else "" for c in row]
                    if any(cells): rows.append(" | ".join(cells))
            return "\n".join(rows)
        except: return ""

    def _vision_page(self, page, pn):
        if not self.groq_client: return ""
        try:
            img = page.to_image(resolution=100).original
            return self._vision(img, f"page {pn} of a scanned document. Extract ALL text — names, numbers, dates, IDs, amounts.")
        except Exception as e:
            logger.warning(f"Vision page {pn} failed: {e}"); return ""

    def _page_images(self, page, pn):
        if not self.groq_client: return []
        out = []
        try:
            from PIL import Image
            for obj in page.images[:2]:
                try:
                    data = obj.get("stream",{})
                    if hasattr(data,"get_data"): raw=data.get_data()
                    elif isinstance(data,bytes): raw=data
                    else: continue
                    img = Image.open(io.BytesIO(raw))
                    if img.width<80 or img.height<80: continue
                    d = self._vision(img, f"a chart/diagram on page {pn}. Describe all data, values, labels.")
                    if d: out.append({"content":f"[IMAGE p{pn}]: {d}","metadata":{"page":pn}})
                except: continue
        except: pass
        return out

    def _vision(self, pil_img, context):
        try:
            import base64
            from PIL import Image
            if pil_img.mode not in ("RGB","L"): pil_img = pil_img.convert("RGB")
            mx = 1024
            w,h = pil_img.size
            if max(w,h) > mx:
                sc = mx/max(w,h); pil_img = pil_img.resize((int(w*sc),int(h*sc)), Image.LANCZOS)
            buf = io.BytesIO(); pil_img.save(buf, format="PNG", optimize=True)
            b64 = base64.b64encode(buf.getvalue()).decode()
            vm  = getattr(self.config,"GROQ_VISION_MODEL","meta-llama/llama-4-scout-17b-16e-instruct") if self.config else "meta-llama/llama-4-scout-17b-16e-instruct"
            r   = self.groq_client.chat.completions.create(
                model=vm, max_tokens=2000,
                messages=[{"role":"user","content":[
                    {"type":"image_url","image_url":{"url":f"data:image/png;base64,{b64}"}},
                    {"type":"text","text":f"This is {context} Be thorough and complete."}
                ]}]
            )
            result = r.choices[0].message.content.strip()
            logger.info(f"Vision returned {len(result)} chars")
            return result
        except Exception as e:
            logger.error(f"Vision error: {e}"); return ""

    # ── IMAGE file ───────────────────────────────────────
    def _image(self, fp):
        try:
            from PIL import Image
            img  = Image.open(fp)
            desc = self._vision(img, "an uploaded image. Extract all text and describe all content.")
            return [{"content":f"[IMAGE]\n{desc}","metadata":{}}] if desc else []
        except Exception as e:
            logger.error(f"Image load failed: {e}"); return []

    # ── TEXT ─────────────────────────────────────────────
    def _text(self, fp):
        for enc in ["utf-8","utf-8-sig","latin-1","cp1252","utf-16"]:
            try:
                with open(fp,"r",encoding=enc,errors="replace") as f:
                    return [{"content":f.read(),"metadata":{"encoding":enc}}]
            except: continue
        with open(fp,"rb") as f:
            return [{"content":f.read().decode("utf-8",errors="replace"),"metadata":{}}]

    # ── DOCX ─────────────────────────────────────────────
    def _docx(self, fp):
        try:
            from docx import Document
            doc=Document(fp); sections,cur=[],[]
            for para in doc.paragraphs:
                t=para.text.strip()
                if not t: continue
                if para.style.name.startswith("Heading") and cur:
                    sections.append({"content":"\n".join(cur),"metadata":{}}); cur=[t]
                else: cur.append(t)
            for table in doc.tables:
                rows=[]
                for row in table.rows:
                    cells=[c.text.strip() for c in row.cells if c.text.strip()]
                    if cells: rows.append(" | ".join(cells))
                if rows: sections.append({"content":"TABLE:\n"+"\n".join(rows),"metadata":{"type":"table"}})
            if cur: sections.append({"content":"\n".join(cur),"metadata":{}})
            return sections or [{"content":"\n".join(p.text for p in doc.paragraphs if p.text.strip()),"metadata":{}}]
        except Exception as e:
            logger.error(f"DOCX failed: {e}"); return []

    # ── CSV ──────────────────────────────────────────────
    def _csv(self, fp):
        rows=[]
        for enc in ["utf-8-sig","utf-8","latin-1","cp1252"]:
            try:
                with open(fp,"r",encoding=enc,errors="replace") as f:
                    rows=list(csv.DictReader(f)); break
            except: continue
        chunks=[]
        for i in range(0,len(rows),50):
            b=rows[i:i+50]
            chunks.append({"content":"\n".join(", ".join(f"{k}: {v}" for k,v in r.items()) for r in b),"metadata":{"rows":f"{i+1}-{i+len(b)}"}})
        return chunks

    # ── XLSX ─────────────────────────────────────────────
    def _xlsx(self, fp):
        try:
            import openpyxl; wb=openpyxl.load_workbook(fp,read_only=True,data_only=True); chunks=[]
            for sh in wb.sheetnames:
                ws=wb[sh]; rows=[]; hdrs=[]
                for i,row in enumerate(ws.iter_rows(values_only=True)):
                    if i==0: hdrs=[str(c) if c is not None else f"c{j}" for j,c in enumerate(row)]
                    else:
                        d={hdrs[j]:str(v) if v is not None else "" for j,v in enumerate(row) if j<len(hdrs)}
                        rows.append(d)
                    if i>10000: break
                for i in range(0,len(rows),50):
                    b=rows[i:i+50]
                    chunks.append({"content":f"Sheet:{sh}\n"+"\n".join(", ".join(f"{k}:{v}" for k,v in r.items()) for r in b),"metadata":{"sheet":sh}})
            wb.close(); return chunks
        except Exception as e:
            logger.error(f"XLSX: {e}"); return []

    # ── PPTX ─────────────────────────────────────────────
    def _pptx(self, fp):
        try:
            from pptx import Presentation; prs=Presentation(fp); slides=[]
            for i,slide in enumerate(prs.slides):
                texts=[]
                for shape in slide.shapes:
                    if hasattr(shape,"text") and shape.text.strip(): texts.append(shape.text.strip())
                    if shape.has_table:
                        rows=[]
                        for row in shape.table.rows:
                            cells=[c.text.strip() for c in row.cells if c.text.strip()]
                            if cells: rows.append(" | ".join(cells))
                        if rows: texts.append("TABLE:\n"+"\n".join(rows))
                if texts: slides.append({"content":"\n".join(texts),"metadata":{"slide":i+1}})
            return slides
        except Exception as e:
            logger.error(f"PPTX: {e}"); return []

    # ── JSON / XML / HTML ────────────────────────────────
    def _json(self, fp):
        try:
            with open(fp,"r",encoding="utf-8",errors="replace") as f: data=json.load(f)
            text=json.dumps(data,indent=2,ensure_ascii=False)
            return [{"content":text[i:i+3000],"metadata":{"part":i//3000+1}} for i in range(0,len(text),3000)]
        except Exception as e:
            logger.error(f"JSON: {e}"); return []

    def _xml(self, fp):
        try:
            tree=ET.parse(fp)
            texts=[f"{e.tag}: {e.text.strip()}" for e in tree.getroot().iter() if e.text and e.text.strip()]
            return [{"content":"\n".join(texts),"metadata":{}}]
        except Exception as e:
            logger.error(f"XML: {e}"); return []

    def _html(self, fp):
        try:
            from bs4 import BeautifulSoup
            with open(fp,"r",encoding="utf-8",errors="replace") as f: soup=BeautifulSoup(f,"html.parser")
            for t in soup(["script","style","nav","footer"]): t.decompose()
            text="\n".join(l.strip() for l in soup.get_text("\n").splitlines() if l.strip())
            return [{"content":text,"metadata":{}}]
        except Exception as e:
            logger.error(f"HTML: {e}"); return []

    # ── HELPERS ──────────────────────────────────────────
    def _rm_noise(self, pages):
        if len(pages)<4: return pages
        all_lines=[]
        for p in pages: all_lines.extend(set(l.strip() for l in p.get("content","").splitlines() if l.strip()))
        freq=Counter(all_lines); thresh=max(3,len(pages)*0.4)
        noise={l for l,c in freq.items() if c>=thresh and len(l)<120}
        if not noise: return pages
        for p in pages:
            p["content"]="\n".join(l for l in p.get("content","").splitlines() if l.strip() not in noise)
        return pages

    def _mask_pii(self, pages):
        pats=[
            (re.compile(r'\b[A-Z]{5}[0-9]{4}[A-Z]\b'),                           "[PAN-REDACTED]"),
            (re.compile(r'\b\d{4}[\s\-]?\d{4}[\s\-]?\d{4}\b'),                   "[AADHAAR-REDACTED]"),
            (re.compile(r'(\+91[\s\-]?|0)?[6-9]\d{9}\b'),                         "[PHONE-REDACTED]"),
            (re.compile(r'\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b'),"[EMAIL-REDACTED]"),
            (re.compile(r'\b[A-Z]{4}0[A-Z0-9]{6}\b'),                             "[IFSC-REDACTED]"),
        ]
        for p in pages:
            t=p.get("content","")
            for pat,rep in pats: t=pat.sub(rep,t)
            p["content"]=t
        return pages
